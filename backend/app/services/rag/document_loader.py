from __future__ import annotations

import asyncio
from pathlib import Path


class DocumentLoader:
    """多格式文档加载器（教学版）。

    原则：
    - 单一职责：只负责“把文件转成文本”；
    - 统一接口：load_text(path, ext)；
    - 异常吞吐：返回空串或提示，不让上传流程直接崩掉。
    """

    async def load_text(self, path: Path, ext: str) -> str:
        # asyncio.to_thread: 把阻塞型库调用放到线程，避免卡住事件循环。
        if ext in {'txt', 'md'}:
            return await asyncio.to_thread(path.read_text, encoding='utf-8', errors='ignore')
        if ext == 'pdf':
            return await asyncio.to_thread(self._load_pdf, path)
        if ext == 'pptx':
            return await asyncio.to_thread(self._load_pptx, path)
        if ext == 'docx':
            return await asyncio.to_thread(self._load_docx, path)
        if ext in {'png', 'jpg', 'jpeg'}:
            return await asyncio.to_thread(self._load_image_ocr, path)
        return ''

    def _load_pdf(self, path: Path) -> str:
        """PDF 文本提取。扫描版通常会返回空文本，需要 OCR。"""
        try:
            from pypdf import PdfReader
        except ImportError:
            return 'PDF解析依赖缺失：请安装 pypdf。'
        try:
            reader = PdfReader(str(path))
            return '\n'.join([(p.extract_text() or '').strip() for p in reader.pages if (p.extract_text() or '').strip()])
        except Exception:
            return ''

    def _load_pptx(self, path: Path) -> str:
        """PPTX 文本提取：遍历 slide shapes 收集可见文本。"""
        try:
            from pptx import Presentation
        except ImportError:
            return 'PPTX解析依赖缺失：请安装 python-pptx。'
        try:
            prs = Presentation(str(path))
            out: list[str] = []
            for s in prs.slides:
                for sh in s.shapes:
                    if hasattr(sh, 'text'):
                        txt = (sh.text or '').strip()
                        if txt:
                            out.append(txt)
            return '\n'.join(out)
        except Exception:
            return ''

    def _load_docx(self, path: Path) -> str:
        """DOCX 文本提取：读取段落文本。"""
        try:
            from docx import Document
        except ImportError:
            return 'DOCX解析依赖缺失：请安装 python-docx。'
        try:
            doc = Document(str(path))
            return '\n'.join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
        except Exception:
            return ''

    def _load_image_ocr(self, path: Path) -> str:
        """图片 OCR：依赖 pytesseract + 本机 tesseract。"""
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            return '图片OCR依赖缺失：请安装 pytesseract 和 pillow。'
        try:
            return pytesseract.image_to_string(Image.open(path), lang='chi_sim+eng').strip()
        except Exception:
            return ''
